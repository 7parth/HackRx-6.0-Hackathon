from fastapi import APIRouter, HTTPException, Header
from typing import List, Dict, Any
import requests
import tempfile
import os
import logging
import re
from docx import Document as DocxDocument
from pydantic import BaseModel, Field
from ..RAG.rag_llm import AdaptiveGeneralLLMDocumentQASystem
import fitz
import time
from concurrent.futures import ThreadPoolExecutor
import openpyxl
from pptx import Presentation
import pytesseract
from PIL import Image
import io
import PyPDF2
import re
from typing import Dict, List
import logging
import docx
import io
import openpyxl
from pptx import Presentation
import pytesseract
from PIL import Image
from collections import OrderedDict
import hashlib


logger = logging.getLogger(__name__)
router = APIRouter(tags=["hackrx"], prefix="/api/v1")


DOCUMENT_CACHE = OrderedDict()
MAX_CACHE_SIZE = 50 
DOCUMENT_CACHE = {}  # Global cache for known documents
KNOWN_DOCUMENTS = [
    "https://example.com/policy1.pdf",
    "https://example.com/terms.docx",
    # Add more known document URLs
]


def preload_known_documents():
    """Pre-cache known documents at startup"""
    downloader = EnhancedDocumentDownloader()
    extractor = FileMetadataExtractor()
    
    for url in KNOWN_DOCUMENTS:
        try:
            logger.info(f"Pre-caching known document: {url}")
            temp_path, _ = downloader.download_from_url(url)
            metadata = extractor.extract_metadata(temp_path)
            file_type = metadata['file_type']
            
            # Extract text based on file type
            if file_type == 'pdf':
                text = downloader.extract_text_from_pdf_enhanced(temp_path)
            elif file_type in ['docx', 'doc']:
                text = downloader.extract_text_from_docx(temp_path)
            elif file_type in ['xlsx', 'xls']:
                text = downloader.extract_text_from_xlsx(temp_path)
            elif file_type in ['pptx', 'ppt']:
                text = downloader.extract_text_from_pptx(temp_path)
            elif file_type in ['jpg', 'jpeg', 'png', 'gif']:
                text = downloader.extract_text_from_image(temp_path)
            else:
                # Fallback to text extraction
                with open(temp_path, 'r', encoding='utf-8') as f:
                    text = f.read()
            
            # Generate cache key
            cache_key = hashlib.sha256(url.encode()).hexdigest()
            
            # Store in cache
            DOCUMENT_CACHE[cache_key] = {
                'text': text,
                'metadata': metadata,
                'url': url
            }
            
            os.unlink(temp_path)
            logger.info(f"Cached document: {url} ({len(text)} chars)")
        except Exception as e:
            logger.error(f"Failed to cache {url}: {str(e)}")

# Call during module initialization
preload_known_documents()

class HackRXRequest(BaseModel):
    documents: str = Field(..., description="Document URL or plain text content")
    questions: List[str] = Field(..., min_items=1, max_items=50, description="List of questions (max 50)") # type: ignore


class HackRXResponse(BaseModel):
    answers: List[str] = Field(..., description="Answers corresponding to the questions")


hackrx_rag_system = None


def get_hackrx_rag_system():
    """Optimized singleton pattern for RAG system"""
    global hackrx_rag_system
    if hackrx_rag_system is None:
        hackrx_rag_system = AdaptiveGeneralLLMDocumentQASystem("")
    return hackrx_rag_system


class FileMetadataExtractor:
    def __init__(self):
        # Ordered from MOST specific to LEAST specific
        self.signatures = [
            (b'%PDF-', 'pdf'),
            (b'PK\x03\x04\x14\x00\x06\x00', 'pptx'),  # PPTX specific
            (b'PK\x03\x04\x14\x00\x08\x00', 'xlsx'),  # XLSX specific
            (b'PK\x03\x04\x14\x00\x00\x00', 'docx'),  # DOCX specific
            (b'PK\x03\x04', 'docx'),  # Generic Office fallback
            (b'\xFF\xD8\xFF', 'jpg'),
            (b'\x89PNG\r\n\x1a\n', 'png'),
            (b'GIF87a', 'gif'),
            (b'GIF89a', 'gif')
        ]
    
    def get_file_type(self, path: str) -> str:
        """Detect file type using magic bytes with content-type fallback"""
        try:
            with open(path, 'rb') as f:
                header = f.read(32)
            
            # Check specific signatures first
            for sig, ftype in self.signatures:
                if header.startswith(sig):
                    return ftype
                
            # Fallback to extension
            ext = os.path.splitext(path)[1].lower()
            if ext:
                return ext[1:]  # Remove dot
            
            return 'unknown'
        except Exception as e:
            logger.error(f"File type detection failed: {str(e)}")
            return 'unknown'
    
    def extract_metadata(self, path: str) -> Dict[str, Any]:
        """Extract metadata based on detected file type"""
        file_type = self.get_file_type(path)
        metadata = {'file_type': file_type, 'page_count': 0}
        
        try:
            if file_type == 'pdf':
                return self.extract_pdf_metadata(path) | metadata
            elif file_type in ['docx', 'doc']:
                return self.extract_docx_metadata(path) | metadata
            elif file_type in ['xlsx', 'xls']:
                return self.extract_xlsx_metadata(path) | metadata
            elif file_type in ['pptx', 'ppt']:
                return self.extract_pptx_metadata(path) | metadata
            elif file_type in ['jpg', 'jpeg', 'png', 'gif']:
                return self.extract_image_metadata(path) | metadata
            else:
                logger.warning(f"Unsupported file type: {file_type}")
                return metadata
        except Exception as e:
            logger.error(f"Metadata extraction failed: {str(e)}")
            return metadata

    def extract_pdf_metadata(self, path: str) -> Dict[str, Any]:
        file_obj = None
        try:
            file_obj = open(path, 'rb')
            pdf_reader = PyPDF2.PdfReader(file_obj)
            
            metadata = pdf_reader.metadata or {}
            first_page_text = ""
            if len(pdf_reader.pages) > 0:
                first_page_text = pdf_reader.pages[0].extract_text() or ""
                first_page_text = first_page_text[:1000]
            
            return {
                'title': str(metadata.get('/Title', '')).lower(),
                'subject': str(metadata.get('/Subject', '')).lower(),
                'keywords': str(metadata.get('/Keywords', '')).lower(),
                'creator': str(metadata.get('/Creator', '')).lower(),
                'producer': str(metadata.get('/Producer', '')).lower(),
                'author': str(metadata.get('/Author', '')).lower(),
                'first_page_sample': first_page_text.lower()[:500],
                'page_count': len(pdf_reader.pages),
            }
        except Exception as e:
            logging.error(f"PDF metadata error: {e}")
            return {}
        finally:
            if file_obj:
                file_obj.close()

    def extract_docx_metadata(self, path: str) -> Dict[str, Any]:
        try:
            doc = docx.Document(path)
            first_content = ""
            for paragraph in doc.paragraphs[:5]:
                first_content += paragraph.text + " "
                if len(first_content) > 500:
                    break
            
            core_props = doc.core_properties
            return {
                'title': str(core_props.title or '').lower(),
                'subject': str(core_props.subject or '').lower(),
                'keywords': str(core_props.keywords or '').lower(),
                'creator': str(core_props.author or '').lower(),
                'author': str(core_props.author or '').lower(),
                'first_page_sample': first_content.lower()[:500],
                'page_count': len(doc.paragraphs),
            }
        except Exception as e:
            logging.error(f"DOCX metadata error: {e}")
            return {}

    def extract_xlsx_metadata(self, path: str) -> Dict[str, Any]:
        try:
            wb = openpyxl.load_workbook(path)
            first_page_sample = ""
            sheet = wb.active
            for row in sheet.iter_rows(max_row=10, values_only=True):
                first_page_sample += " ".join(str(cell) for cell in row if cell) + "\n"
            
            return {
                'title': wb.properties.title.lower() if wb.properties.title else "",
                'first_page_sample': first_page_sample.lower()[:500],
                'sheet_count': len(wb.sheetnames),
            }
        except Exception as e:
            logging.error(f"XLSX metadata error: {e}")
            return {}

    def extract_pptx_metadata(self, path: str) -> Dict[str, Any]:
        try:
            prs = Presentation(path)
            first_page_sample = ""
            if len(prs.slides) > 0:
                slide = prs.slides[0]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                first_page_sample += run.text + " "
                    elif shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            for cell in row.cells:
                                if cell.text_frame:
                                    first_page_sample += cell.text_frame.text + " | "
            
            return {
                'title': prs.core_properties.title.lower() if prs.core_properties.title else "",
                'first_page_sample': first_page_sample.lower()[:500],
                'slide_count': len(prs.slides),
            }
        except Exception as e:
            logging.error(f"PPTX metadata error: {e}")
            return {}

    def extract_image_metadata(self, path: str) -> Dict[str, Any]:
        try:
            img = Image.open(path)
            return {
                'title': "",
                'first_page_sample': "Image content detected",
            }
        except Exception as e:
            logging.error(f"Image metadata error: {e}")
            return {}

class EnhancedDocumentDownloader:
    def __init__(self):
        self.max_file_size = 1024 * 1024 * 1024  # 1 GB
        self.page_number_pattern = re.compile(r'\n\s*\d+\s*\n')
        self.page_header_pattern = re.compile(r'\n\s*Page\s+\d+.*?\n', re.IGNORECASE)
        self.whitespace_pattern = re.compile(r'\s+')
        self.line_break_pattern = re.compile(r'\n\s*\n\s*\n')
            # Add document type mapping
        
    def download_from_url(self, url: str) -> tuple[str, str]:
        """Download document from URL with enhanced error handling"""
        try:
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            }
            response = requests.get(url, timeout=30, stream=True, headers=headers)
            response.raise_for_status()

            content_length = response.headers.get('content-length')
            if content_length and int(content_length) > self.max_file_size:
                raise HTTPException(status_code=413, detail=f"File too large (max {self.max_file_size // (1024 * 1024)}MB)")

            # Create temporary file with no extension
            temp_file = tempfile.NamedTemporaryFile(delete=False)
            total_size = 0
            
            # Stream download content
            for chunk in response.iter_content(chunk_size=16384):
                if chunk:
                    total_size += len(chunk)
                    if total_size > self.max_file_size:
                        temp_file.close()
                        os.unlink(temp_file.name)
                        raise HTTPException(status_code=413, detail=f"File too large (max {self.max_file_size // (1024 * 1024)}MB)")
                    temp_file.write(chunk)
            
            temp_file.close()
            return temp_file.name, 'unknown'  # File type will be determined later

        except requests.RequestException as e:
            logger.error(f"Download failed: {str(e)}")
            raise HTTPException(status_code=400, detail=f"Failed to download document: {str(e)}")

    def extract_text_from_pdf_enhanced(self, path: str) -> str:
        """Performance-optimized PDF extraction"""
        try:
            doc = fitz.open(path)
            page_count = len(doc)
            logger.info(f"Processing PDF with {page_count} pages")
            max_pages = page_count

            if page_count > 50:
                text_pages = self._extract_pdf_parallel(doc, max_pages)
            else:
                text_pages = self._extract_pdf_sequential(doc, max_pages)

            doc.close()

            combined_text = "\n\n".join(text_pages)
            logger.info(f"Extracted {len(combined_text)} characters from PDF")
            return combined_text

        except Exception as e:
            logger.error(f"PDF extraction failed: {str(e)}")
            raise HTTPException(status_code=400, detail=f"Failed to extract PDF content: {str(e)}")

    def _extract_pdf_parallel(self, doc, max_pages: int) -> List[str]:
        def extract_page(page_num):
            try:
                page = doc[page_num]
                page_text = self._extract_text_with_layout_awareness(page)
                if page_text.strip():
                    return self._clean_text(page_text)
                return ""
            except Exception as e:
                logger.warning(f"Error extracting page {page_num}: {e}")
                return ""

        with ThreadPoolExecutor(max_workers=4) as executor:
            text_pages = list(executor.map(extract_page, range(max_pages)))

        return [page for page in text_pages if page.strip()]

    def _extract_pdf_sequential(self, doc, max_pages: int) -> List[str]:
        text_pages = []
        for page_num in range(max_pages):
            page = doc[page_num]
            page_text = self._extract_text_with_layout_awareness(page)
            if page_text.strip():
                cleaned_text = self._clean_text(page_text)
                text_pages.append(cleaned_text)
        return text_pages

    def extract_text_from_docx(self, path: str) -> str:
        try:
            doc = DocxDocument(path)
            paragraphs = []

            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:
                    paragraphs.append(text)

            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        paragraphs.append(" | ".join(row_text))

            combined_text = "\n\n".join(paragraphs)
            logger.info(f"Extracted {len(combined_text)} characters from DOCX")
            return combined_text

        except Exception as e:
            logger.error(f"DOCX extraction failed: {str(e)}")
            raise HTTPException(status_code=400, detail=f"Failed to extract DOCX content: {str(e)}")
    
    def extract_text_from_xlsx(self, path: str) -> str:
        """Extract text from Excel files"""
        try:
            wb = openpyxl.load_workbook(path)
            text_parts = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"\n\n--- Sheet: {sheet_name} ---\n")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) for cell in row if cell]
                    if row_text:
                        text_parts.append(" | ".join(row_text))
            
            combined_text = "\n".join(text_parts)
            logger.info(f"Extracted {len(combined_text)} characters from XLSX")
            return combined_text
        except Exception as e:
            logger.error(f"XLSX extraction failed: {str(e)}")
            raise HTTPException(status_code=400, detail=f"Failed to extract XLSX content: {str(e)}")
    
    def extract_text_from_pptx(self, path: str) -> str:
        """Extract text from PowerPoint files"""
        try:
            prs = Presentation(path)
            text_parts = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = f"\n\n--- Slide {i+1} ---\n"
                for shape in slide.shapes:
                    # Handle text frames
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                slide_text += run.text + " "
                            slide_text += "\n"
                    # Handle tables
                    elif shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text_frame:
                                    row_text.append(cell.text_frame.text.strip())
                            if any(row_text):
                                slide_text += " | ".join(row_text) + "\n"
                    # Handle group shapes (nested shapes)
                    elif shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                        for sub_shape in shape.shapes:
                            if sub_shape.has_text_frame:
                                for paragraph in sub_shape.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        slide_text += run.text + " "
                text_parts.append(slide_text)
            
            combined_text = "\n".join(text_parts)
            logger.info(f"Extracted {len(combined_text)} characters from PPTX")
            return combined_text
        except Exception as e:
            logger.error(f"PPTX extraction failed: {str(e)}")
            raise HTTPException(status_code=400, detail=f"Failed to extract PPTX content: {str(e)}")
    
    def extract_text_from_image(self, path: str) -> str:
        """Extract text from images using OCR"""
        try:
            img = Image.open(path)
            text = pytesseract.image_to_string(img)
            logger.info(f"Extracted {len(text)} characters from image via OCR")
            return text
        except Exception as e:
            logger.error(f"Image OCR failed: {str(e)}")
            raise HTTPException(status_code=400, detail=f"Failed to extract text from image: {str(e)}")

    def _extract_text_with_layout_awareness(self, page) -> str:
        try:
            blocks = page.get_text("dict")
            if not blocks.get("blocks"):
                return page.get_text()
            page_width = page.rect.width
            left_blocks = []
            right_blocks = []
            for block in blocks["blocks"]:
                if "lines" not in block:
                    continue

                block_bbox = block["bbox"]
                block_center_x = (block_bbox[0] + block_bbox[2]) / 2
                block_text = ""

                for line in block["lines"]:
                    line_text = ' '.join([span.get("text", "") for span in line.get("spans", [])])
                    if line_text.strip():
                        block_text += line_text + " "

                if block_text.strip():
                    if block_center_x < page_width * 0.55:
                        left_blocks.append((block_bbox[1], block_text.strip()))
                    else:
                        right_blocks.append((block_bbox[1], block_text.strip()))

            if len(left_blocks) > 2 and len(right_blocks) > 2:
                left_blocks.sort()
                right_blocks.sort()
                combined_text = []
                combined_text.extend([text for _, text in left_blocks])
                combined_text.extend([text for _, text in right_blocks])
                return "\n\n".join(combined_text)
            else:
                return page.get_text()

        except Exception:
            return page.get_text()

    def _clean_text(self, text: str) -> str:
        text = self.line_break_pattern.sub('\n\n', text)
        text = self.whitespace_pattern.sub(' ', text)
        text = self.page_number_pattern.sub('\n', text)
        text = self.page_header_pattern.sub('\n', text)
        return text.strip()


class AdaptiveParameterSelector:
    def __init__(self):
        self.technical_pattern = re.compile(r'\b\d+\.\d+\b|\b[A-Z]{2,}\b|\([^)]*\)|\b(?:Figure|Table|Section)\s+\d+')

    def get_optimal_parameters(self, document_text: str) -> Dict[str, Any]:
        word_count = len(document_text.split())
        page_estimate = max(1, word_count // 250)
        technical_matches = len(self.technical_pattern.findall(document_text))
        has_technical_content = technical_matches > (word_count * 0.015)

        if page_estimate <= 10:
            chunk_size, chunk_overlap, retriever_k = 1200, 200, 6
        elif page_estimate <= 35:
            chunk_size, chunk_overlap, retriever_k = 1800, 300, 8
        elif page_estimate <= 100:
            chunk_size, chunk_overlap, retriever_k = 2500, 400, 10
        else:
            chunk_size, chunk_overlap, retriever_k = 3000, 500, 12

        if has_technical_content:
            chunk_size = int(chunk_size * 1.1)

        return {
            'chunk_size': chunk_size,
            'chunk_overlap': chunk_overlap,
            'retriever_k': retriever_k,
            'estimated_pages': page_estimate,
            'has_technical_content': has_technical_content
        }


@router.post("/hackrx/run", response_model=HackRXResponse)
def process_document_questions(
    request: HackRXRequest,
    authorization: str = Header(..., description="Bearer token for authentication")
    ):
    start_time = time.time()
    expected_token = "Bearer a087324753b37209904afffffa5ad45b8aac7912c74f61420e7e054237778a95"
    document_text = ""
    metadata = {}
    cache_key = None
    
    # Check if document is in cache
    if request.documents.startswith(("http://", "https://")):
        cache_key = hashlib.sha256(request.documents.encode()).hexdigest()
        
        if cache_key in DOCUMENT_CACHE:
            logger.info(f"Using cached document: {request.documents}")
            cached = DOCUMENT_CACHE[cache_key]
            document_text = cached['text']
            metadata = cached['metadata']
            
            # Move to end to mark as recently used (LRU)
            DOCUMENT_CACHE.move_to_end(cache_key)

    
    # Authorization check
    if authorization != expected_token:
        logger.warning("401 - Invalid authorization token received.")
        raise HTTPException(status_code=401, detail="Invalid authorization token")

    # Field presence checks
    if request.documents is None and request.questions is None:
        logger.warning("400 - Both 'documents' and 'questions' fields are missing.")
        raise HTTPException(status_code=400, detail="Both documents and questions are required")

    if not request.documents:
        logger.warning("400 - 'documents' field is missing or empty.")
        raise HTTPException(status_code=400, detail="'documents' field is required")

    if not request.questions:
        logger.warning("400 - 'questions' field is missing or empty.")
        raise HTTPException(status_code=400, detail="'questions' field is required")

    if not isinstance(request.questions, list):
        logger.warning(f"400 - 'questions' should be a list, got: {type(request.questions)}")
        raise HTTPException(status_code=400, detail="'questions' should be a list of strings")

    if not all(isinstance(q, str) for q in request.questions):
        logger.warning("400 - One or more entries in 'questions' are not strings.")
        raise HTTPException(status_code=400, detail="All questions must be strings.")

    if len(request.questions) > 50:
        logger.warning(f"400 - Too many questions: {len(request.questions)} (max 50 allowed)")
        raise HTTPException(status_code=400, detail="Maximum 50 questions allowed per request")

    # Enhanced logging of request
    logger.info(f"Received request with document source: {'URL' if request.documents.startswith(('http://', 'https://')) else 'TEXT INPUT'}")
    logger.info(f"Document source: {request.documents[:500] + ('...' if len(request.documents) > 500 else '')}")
    logger.info(f"Questions: {request.questions}")

    temp_path = None
    document_text = ""

    try:
        # Process URL-based documents
        if request.documents.startswith(("http://", "https://")):
            downloader = EnhancedDocumentDownloader()
            temp_path, _ = downloader.download_from_url(request.documents)
            
            # Create metadata extractor
            metadata_extractor = FileMetadataExtractor()
            metadata = metadata_extractor.extract_metadata(temp_path)
            file_type = metadata['file_type']
            
            # Log metadata for debugging
            logger.info(f"Document metadata: {metadata}")
            
            # Extract document title or use default
            doc_name = metadata.get('title', 'the document') or "unnamed document"
            if not doc_name or doc_name == "unnamed document":
                # Try to get filename from URL
                if '/' in request.documents:
                    doc_name = request.documents.rsplit('/', 1)[-1].split('?')[0]
            
            # PAGE COUNT CHECK (keep only this)
            page_count = metadata.get('page_count', 0)
            if page_count > 450:
                error_msg = (f"The provided document '{doc_name}' is too long ({page_count} pages). "
                             "We currently do not support documents over 450 pages.")
                return HackRXResponse(answers=[error_msg for _ in request.questions])
            
            # Extract text based on detected file type
            logger.info(f"Extracting text from {file_type} document")
            if file_type == 'pdf':
                document_text = downloader.extract_text_from_pdf_enhanced(temp_path)
            elif file_type in ['docx', 'doc']:
                document_text = downloader.extract_text_from_docx(temp_path)
            elif file_type in ['xlsx', 'xls']:
                document_text = downloader.extract_text_from_xlsx(temp_path)
            elif file_type in ['pptx', 'ppt']:
                document_text = downloader.extract_text_from_pptx(temp_path)
            elif file_type in ['jpg', 'jpeg', 'png', 'gif']:
                document_text = downloader.extract_text_from_image(temp_path)
            else:
                # Try text extraction as fallback
                try:
                    with open(temp_path, 'r', encoding='utf-8') as f:
                        document_text = f.read()
                except UnicodeDecodeError:
                    try:
                        with open(temp_path, 'r', encoding='latin-1') as f:
                            document_text = f.read()
                    except Exception as e:
                        logger.error(f"Text extraction failed: {str(e)}")
                        raise HTTPException(status_code=400, detail="Failed to extract text from document")
                except Exception as e:
                    logger.error(f"Text extraction failed: {str(e)}")
                    raise HTTPException(status_code=400, detail="Failed to extract text from document")
        
        # Process text-based input
        else:
            logger.info("Processing text input document")
            document_text = request.documents.strip()
            # Log first 500 characters of text input
            logger.info(f"Text input sample: {document_text[:500]}{'...' if len(document_text) > 500 else ''}")

        # Validate extracted text
        if not document_text or len(document_text.strip()) < 50:
            raise HTTPException(status_code=400, detail="Document content too short or empty")
        
        logger.info(f"Document contains {len(document_text)} characters")

        parameter_selector = AdaptiveParameterSelector()
        optimal_params = parameter_selector.get_optimal_parameters(document_text)
        logger.info(f"Using performance-optimized parameters: {optimal_params}")

        # Initialize and configure RAG system
        logger.info("Initializing RAG system")
        llm = AdaptiveGeneralLLMDocumentQASystem(
            google_api_key="",
            llm_model="gemini-2.0-flash",
            llm_temperature=0.1,
            llm_max_tokens=400,
            top_p=0.95
        )

        # Prepare document data
        doc_data = [{
            "content": document_text,
            "filename": "hackrx_document.txt",
            "doc_id": "hackrx_doc_1"
        }]

        # Process documents and questions
        logger.info("Loading documents into RAG system")
        llm.load_documents_from_content_adaptive(doc_data)
        
        logger.info(f"Processing {len(request.questions)} questions")
        doc_url = request.documents.strip()

        # Prepare doc_data for vectorizing
        doc_data = [{
            "content": document_text,
            "filename": metadata.get("title", "document.txt") or "document.txt",
            "doc_id": "cached_doc_1"
        }]

        # Try loading from cache
        if llm.try_load_from_url_cache(doc_url):
            logger.info("RAG system loaded from cache.")
        else:
            logger.info("Processing and caching document.")
            llm.load_documents_from_content_adaptive(doc_data)
            llm.save_to_url_cache(doc_url, document_text, metadata)
        answers = llm.process_questions_batch(request.questions)

        # Log processing time
        processing_time = round(time.time() - start_time, 2)
        logger.info(f"Successfully processed request in {processing_time} seconds")
        
        # Log first 3 answers as sample
        sample_answers = answers[:3]
        if len(answers) > 3:
            sample_answers.append("...")
        logger.info(f"Sample answers: {sample_answers}")

        return HackRXResponse(answers=answers)

    except HTTPException:
        raise

    except Exception as e:
        logger.error(f"Internal server error: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Internal server error: {str(e)}")

    finally:
        if temp_path and os.path.exists(temp_path):
            try:
                os.unlink(temp_path)
                logger.debug(f"Cleaned up temporary file: {temp_path}")
            except Exception as e:
                logger.warning(f"Failed to cleanup temp file: {e}")