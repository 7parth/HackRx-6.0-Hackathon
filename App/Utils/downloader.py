import os
import requests
import tempfile
from fastapi import HTTPException
import requests
import tempfile
from typing import List
import fitz
from concurrent.futures import ThreadPoolExecutor
import pytesseract
import pytesseract
import re
import logging
import openpyxl
from pptx import Presentation
from PIL import Image
from docx import Document as DocxDocument

logger = logging.getLogger(__name__)
class DocumentDownloader:
    def download_from_url(self, url: str) -> tuple[str, str]:
        """
        Download file from URL to a temporary location. Returns (path, extension)
        """
        response = requests.get(url)
        response.raise_for_status()

        content_type = response.headers.get("Content-Type", "").lower()
        ext = ""
        if "pdf" in content_type:
            ext = ".pdf"
        elif "word" in content_type:
            ext = ".docx"
        elif "presentationml" in content_type:
            ext = ".pptx"
        elif "spreadsheetml" in content_type:
            ext = ".xlsx"
        elif "image" in content_type:
            ext = ".png"

        _, temp_path = tempfile.mkstemp(suffix=ext)
        with open(temp_path, "wb") as f:
            f.write(response.content)

        return temp_path, ext
    
class EnhancedDocumentDownloader:
    def __init__(self):
        self.max_file_size = 1024 * 1024 * 1024  # 1 GB
        self.page_number_pattern = re.compile(r'\n\s*\d+\s*\n')
        self.page_header_pattern = re.compile(r'\n\s*Page\s+\d+.*?\n', re.IGNORECASE)
        self.whitespace_pattern = re.compile(r'\s+')
        self.line_break_pattern = re.compile(r'\n\s*\n\s*\n')
            # Add document type mapping
        
    def download_from_url(self, url: str) -> tuple[str, str]:

        try:
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            }
            response = requests.get(url, timeout=30, stream=True, headers=headers)
            response.raise_for_status()
            content_length = response.headers.get('content-length')
            if content_length and int(content_length) > self.max_file_size:
                raise HTTPException(status_code=413, detail="File too large (1GB max)")
            # get extension from URL or fallback
            ext = os.path.splitext(url.split('?')[0])[1]
            if ext.lower() not in ('.xlsx', '.xlsm', '.xltx', '.xltm', '.pdf', '.docx', '.pptx', '.png', '.jpg', '.jpeg', '.gif'):
                ext = ''
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as temp_file:
                total_size = 0
                for chunk in response.iter_content(chunk_size=16384):
                    if chunk:
                        total_size += len(chunk)
                        if total_size > self.max_file_size:
                            temp_file.close()
                            os.unlink(temp_file.name)
                            raise HTTPException(status_code=413, detail="File too large (1GB max)")
                        temp_file.write(chunk)
                temp_path = temp_file.name
            # Defensive: Check file is not empty or tiny
            if os.path.getsize(temp_path) < 1024:
                os.unlink(temp_path)
                raise HTTPException(status_code=400, detail="Downloaded file appears truncated/corrupt.")
            return temp_path, 'unknown'
        except requests.RequestException as e:
            logger.error(f"Download failed: {str(e)}")
            raise HTTPException(status_code=400, detail=f"Failed to download document: {str(e)}")
        
    def extract_text_from_pdf_enhanced(self, path: str) -> str:
        """Performance-optimized PDF extraction"""
        try:
            doc = fitz.open(path)
            page_count = len(doc)
            logger.info(f"Processing PDF with {page_count} pages")
            max_pages = page_count

            if page_count > 50:
                text_pages = self._extract_pdf_parallel(doc, max_pages)
            else:
                text_pages = self._extract_pdf_sequential(doc, max_pages)

            doc.close()

            combined_text = "\n\n".join(text_pages)
            logger.info(f"Extracted {len(combined_text)} characters from PDF")
            return combined_text

        except Exception as e:
            logger.error(f"PDF extraction failed: {str(e)}")
            raise HTTPException(status_code=400, detail=f"Failed to extract PDF content: {str(e)}")

    def _extract_pdf_parallel(self, doc, max_pages: int) -> List[str]:
        def extract_page(page_num):
            try:
                page = doc[page_num]
                page_text = self._extract_text_with_layout_awareness(page)
                if page_text.strip():
                    return self._clean_text(page_text)
                return ""
            except Exception as e:
                logger.warning(f"Error extracting page {page_num}: {e}")
                return ""

        with ThreadPoolExecutor(max_workers=4) as executor:
            text_pages = list(executor.map(extract_page, range(max_pages)))

        return [page for page in text_pages if page.strip()]

    def _extract_pdf_sequential(self, doc, max_pages: int) -> List[str]:
        text_pages = []
        for page_num in range(max_pages):
            page = doc[page_num]
            page_text = self._extract_text_with_layout_awareness(page)
            if page_text.strip():
                cleaned_text = self._clean_text(page_text)
                text_pages.append(cleaned_text)
        return text_pages

    def extract_text_from_docx(self, path: str) -> str:
        try:
            doc = DocxDocument(path)
            paragraphs = []

            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:
                    paragraphs.append(text)

            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        paragraphs.append(" | ".join(row_text))

            combined_text = "\n\n".join(paragraphs)
            logger.info(f"Extracted {len(combined_text)} characters from DOCX")
            return combined_text

        except Exception as e:
            logger.error(f"DOCX extraction failed: {str(e)}")
            raise HTTPException(status_code=400, detail=f"Failed to extract DOCX content: {str(e)}")
    
    def extract_text_from_xlsx(self, path: str) -> str:
        """Extracts text from all cells in all sheets of an Excel workbook."""
        import zipfile
        if not zipfile.is_zipfile(path):
            raise HTTPException(status_code=400, detail="Downloaded file is not a valid Excel (.xlsx) file")
        wb = openpyxl.load_workbook(path)
        text_lines = []
        for sheet in wb:
            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else '' for cell in row]
                text_lines.append(' '.join(row_text))
        return '\n'.join(text_lines)
    
    def extract_text_from_pptx(self, path: str) -> str:
        """Comprehensive text extraction from PowerPoint files"""
        try:
            prs = Presentation(path)
            text_parts = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = f"\n\n--- Slide {i+1} ---\n"
                for shape in slide.shapes:
                    # Text frames (standard text boxes)
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                slide_text += run.text + " "
                    
                    # Tables
                    elif shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text_frame:
                                    cell_text = " ".join(
                                        run.text for para in cell.text_frame.paragraphs 
                                        for run in para.runs
                                    )
                                    row_text.append(cell_text.strip())
                            if any(row_text):
                                slide_text += " | ".join(row_text) + "\n"
                    
                    # Grouped shapes
                    elif shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                        for sub_shape in shape.shapes:
                            if sub_shape.has_text_frame:
                                for paragraph in sub_shape.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        slide_text += run.text + " "
                    
                    # Placeholders (title, content, etc)
                    elif shape.is_placeholder:
                        if shape.text:
                            slide_text += shape.text + " "
                
                text_parts.append(slide_text)
            
            combined_text = "\n".join(text_parts)
            logger.info(f"Extracted {len(combined_text)} characters from PPTX")
            return combined_text
        except Exception as e:
            logger.error(f"PPTX extraction failed: {str(e)}")
            raise HTTPException(status_code=400, detail=f"Failed to extract PPTX content: {str(e)}")
    
    def extract_text_from_image(self, path: str) -> str:
        """Robust image text extraction with error handling"""
        try:
            # Validate image file first
            with Image.open(path) as img:
                img.verify()  # Verify file integrity
            
            # Reopen for OCR processing
            with Image.open(path) as img:
                # Convert to RGB if needed (Tesseract requires RGB)
                if img.mode not in ('RGB', 'L'):
                    img = img.convert('RGB')
                
                # Use Tesseract with optimized config
                text = pytesseract.image_to_string(
                    img, 
                    config='--psm 3 --oem 3 -c preserve_interword_spaces=1'
                )
            
            logger.info(f"Extracted {len(text)} characters from image via OCR")
            return text
        except Exception as e:
            logger.error(f"Image OCR failed: {str(e)}")
            raise HTTPException(status_code=400, detail=f"Failed to extract text from image: {str(e)}")

    def _extract_text_with_layout_awareness(self, page) -> str:
        try:
            blocks = page.get_text("dict")
            if not blocks.get("blocks"):
                return page.get_text()
            page_width = page.rect.width
            left_blocks = []
            right_blocks = []
            for block in blocks["blocks"]:
                if "lines" not in block:
                    continue

                block_bbox = block["bbox"]
                block_center_x = (block_bbox[0] + block_bbox[2]) / 2
                block_text = ""

                for line in block["lines"]:
                    line_text = ' '.join([span.get("text", "") for span in line.get("spans", [])])
                    if line_text.strip():
                        block_text += line_text + " "

                if block_text.strip():
                    if block_center_x < page_width * 0.55:
                        left_blocks.append((block_bbox[1], block_text.strip()))
                    else:
                        right_blocks.append((block_bbox[1], block_text.strip()))

            if len(left_blocks) > 2 and len(right_blocks) > 2:
                left_blocks.sort()
                right_blocks.sort()
                combined_text = []
                combined_text.extend([text for _, text in left_blocks])
                combined_text.extend([text for _, text in right_blocks])
                return "\n\n".join(combined_text)
            else:
                return page.get_text()

        except Exception:
            return page.get_text()

    def _clean_text(self, text: str) -> str:
        text = self.line_break_pattern.sub('\n\n', text)
        text = self.whitespace_pattern.sub(' ', text)
        text = self.page_number_pattern.sub('\n', text)
        text = self.page_header_pattern.sub('\n', text)
        return text.strip()